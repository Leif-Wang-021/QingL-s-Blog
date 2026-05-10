from __future__ import annotations

from pathlib import Path
from zipfile import ZipFile
import math
import re
import xml.etree.ElementTree as ET

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(r"C:\Users\26981\Desktop\报告")
PREVIEW_DIR = OUT_DIR / "深度学习汇报_preview"
PPTX_PATH = OUT_DIR / "深度学习汇报.pptx"
MONTAGE_PATH = PREVIEW_DIR / "montage.png"

SLIDE_W, SLIDE_H = 13.333, 7.5
PX_W, PX_H = 1920, 1080

FONT = "HarmonyOS Sans SC"
FONT_REG = Path(r"C:\Windows\Fonts\HarmonyOS_Sans_SC_Regular.ttf")
FONT_BOLD = Path(r"C:\Windows\Fonts\HarmonyOS_Sans_SC_Bold.ttf")

INK = "111827"
MUTED = "667085"
LIGHT = "F8FAFC"
PAPER = "FFFFFF"
LINE = "D0D5DD"
SLATE = "1F2937"
WHITE = "FFFFFF"
CYAN = "0891B2"
BLUE = "2563EB"
GREEN = "059669"
AMBER = "D97706"
PURPLE = "7C3AED"
RED = "DC2626"


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))


def font_px(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    path = FONT_BOLD if bold and FONT_BOLD.exists() else FONT_REG
    if not path.exists():
        path = Path(r"C:\Windows\Fonts\msyh.ttc")
    return ImageFont.truetype(str(path), max(8, int(size * 2)))


def add_text(slide, text, x, y, w, h, size=22, color=INK, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, line_spacing=1.03):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return box


def add_bg(slide, color=LIGHT):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(color)


def add_title(slide, title, subtitle="", dark=False, accent=CYAN):
    title_color = WHITE if dark else INK
    sub_color = "CBD5E1" if dark else MUTED
    add_text(slide, title, 0.68, 0.45, 11.3, 0.62, 28, title_color, True)
    if subtitle:
        add_text(slide, subtitle, 0.70, 1.08, 10.9, 0.34, 13, sub_color)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.70), Inches(1.48), Inches(1.15), Inches(0.045))
    rule.fill.solid()
    rule.fill.fore_color.rgb = rgb(accent)
    rule.line.fill.background()


def add_footer(slide, n: int, dark=False):
    c = "94A3B8" if dark else "98A2B3"
    add_text(slide, "Deep Learning × 通信算法", 0.68, 7.08, 3.1, 0.2, 9, c)
    add_text(slide, f"{n:02d}", 12.35, 7.08, 0.35, 0.2, 9, c, align=PP_ALIGN.RIGHT)


def add_line(slide, x1, y1, x2, y2, color=LINE, width=1.8, arrow=False):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if arrow:
        line.line.end_arrowhead = True
    return line


def add_pill(slide, text, x, y, w, h, fill, color=WHITE, size=13, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.fill.background()
    shape.adjustments[0] = 0.45
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return shape


def add_node(slide, text, x, y, w, h, fill, size=14, color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.fill.background()
    shape.adjustments[0] = 0.18
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = rgb(color)
    return shape


def add_model_slide(prs, blank, n, title, subtitle, accent, core, diagram, fit, diff, comm):
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, title, subtitle, accent=accent)
    add_text(s, "核心理解", 0.85, 2.0, 1.2, 0.28, 13, accent, True)
    add_text(s, core, 0.85, 2.38, 5.1, 0.9, 24, INK, True, line_spacing=1.05)
    x0, y0 = 0.9, 4.15
    for i, (txt, col) in enumerate(diagram):
        x = x0 + i * 1.45
        add_node(s, txt, x, y0, 1.05, 0.45, col, 12)
        if i < len(diagram) - 1:
            add_line(s, x + 1.08, y0 + 0.23, x + 1.38, y0 + 0.23, "94A3B8", 1.4, True)

    groups = [("适合任务", fit, GREEN), ("主要区别", diff, BLUE), ("通信接口", comm, AMBER)]
    for i, (head, body, col) in enumerate(groups):
        y = 2.05 + i * 1.18
        add_pill(s, head, 7.1, y, 1.05, 0.33, col, size=10)
        add_text(s, body, 8.35, y - 0.02, 3.9, 0.56, 16, INK, True, line_spacing=1.08)
    add_footer(s, n)
    return s


def add_table(slide, x, y, colw, rows, header_color=SLATE, row_h=0.46, font_size=12):
    total_w = sum(colw)
    for r, row in enumerate(rows):
        yy = y + r * row_h
        fill = header_color if r == 0 else ("FFFFFF" if r % 2 == 1 else "F1F5F9")
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(yy), Inches(total_w), Inches(row_h))
        rect.fill.solid()
        rect.fill.fore_color.rgb = rgb(fill)
        rect.line.fill.background()
        xx = x
        for c, cell in enumerate(row):
            color = WHITE if r == 0 else INK
            bold = r == 0
            add_text(slide, cell, xx + 0.08, yy + 0.09, colw[c] - 0.12, row_h - 0.1, font_size, color, bold)
            xx += colw[c]


def create_deck():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    slides = []

    # 1
    s = prs.slides.add_slide(blank)
    add_bg(s, SLATE)
    add_text(s, "学习汇报 / Deep Learning × Communication", 0.72, 0.58, 5.1, 0.35, 13, "CBD5E1")
    add_text(s, "Deep Learning\n相关网络原理", 0.72, 1.75, 7.1, 1.4, 42, WHITE, True, line_spacing=0.9)
    add_text(s, "任务适配与通信结合", 0.76, 3.50, 5.8, 0.62, 30, "D1FAE5", True)
    add_text(s, "八类常见网络的学习整理", 0.78, 4.25, 4.6, 0.35, 16, "CBD5E1")
    add_line(s, 8.0, 2.05, 11.6, 2.05, "94A3B8", 1.6)
    for i, (txt, col) in enumerate([("原理", CYAN), ("任务", AMBER), ("通信", GREEN)]):
        add_node(s, txt, 8.05 + i * 1.35, 2.52, 0.9, 0.55, col, 14)
    add_footer(s, 1, dark=True)
    slides.append(("cover", "Deep Learning 相关网络原理", "任务适配与通信结合", True))

    # 2
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "汇报要解决什么", "直接对应学长给的三个问题。", accent=AMBER)
    asks = [
        ("1", "现在 deep learning 相关网络的原理", "八类常见网络的基本原理", CYAN),
        ("2", "各自适合什么任务、有什么区别", "比较适合的数据结构、任务类型和局限", BLUE),
        ("3", "可以跟通信现在哪些任务结合", "对应到信道估计、检测、CSI 反馈、资源分配等问题", GREEN),
    ]
    for i, (num, q, a, col) in enumerate(asks):
        y = 2.05 + i * 1.35
        add_pill(s, num, 0.95, y + 0.02, 0.42, 0.42, col, size=14)
        add_text(s, q, 1.62, y, 4.45, 0.32, 17, INK, True)
        add_line(s, 6.35, y + 0.2, 6.95, y + 0.2, "94A3B8", 1.5, True)
        add_text(s, a, 7.25, y - 0.03, 4.9, 0.5, 18, col, True)
    add_footer(s, 2)
    slides.append(("agenda", "汇报要解决什么", "", False))

    # 3
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "模型分类", "这八类模型可以按任务和数据结构粗略分成几组。", accent=CYAN)
    groups = [
        ("基础前馈类", "FNN", "适合已经整理好的特征输入", BLUE),
        ("空间特征类", "CNN", "适合图像这类规则网格数据", CYAN),
        ("序列建模类", "RNN / LSTM / GRU", "适合文本、语音、时间序列", GREEN),
        ("生成与表示学习类", "GAN / Autoencoder", "前者偏生成，后者偏压缩和重建", AMBER),
        ("注意力与大模型类", "Transformer", "适合全局依赖建模和大规模序列任务", PURPLE),
    ]
    for i, (g, m, d, col) in enumerate(groups):
        y = 2.0 + i * 0.82
        add_text(s, g, 0.95, y, 2.0, 0.28, 15, col, True)
        add_node(s, m, 3.15, y - 0.08, 2.0, 0.42, col, 12)
        add_text(s, d, 5.55, y - 0.01, 5.4, 0.32, 16, INK, True)
    add_footer(s, 3)
    slides.append(("classify", "模型分类", "", False))

    model_specs = [
        (4, "FNN：最基础的输入输出映射", "固定维度输入，学习非线性映射。",
         BLUE, "把一组整理好的特征，逐层变换成分类、回归或决策结果。",
         [("特征", BLUE), ("全连接", BLUE), ("输出", GREEN)],
         "分类、回归、检测、估计。", "不主动利用空间结构或时间顺序。", "信道估计、信号检测、资源分配近似。"),
        (5, "CNN：从局部区域提取特征", "用卷积核在网格数据上滑动。",
         CYAN, "同一个卷积核在不同位置复用，学习局部模式，再逐层组合成更高层特征。",
         [("局部窗口", CYAN), ("卷积核", CYAN), ("特征图", GREEN)],
         "图像、频谱图、信道矩阵、IQ 局部结构。", "比 FNN 更懂相邻位置之间的关系。", "调制识别、频谱感知、CSI 空间特征提取。"),
        (6, "RNN：带记忆的序列模型", "当前状态依赖当前输入和上一状态。",
         GREEN, "模型读到当前位置时，会带着前面信息的摘要继续往后处理。",
         [("x1", GREEN), ("h1", GREEN), ("x2", GREEN), ("h2", AMBER)],
         "文本、语音、时间序列。", "能处理不定长序列，但长序列容易遗忘。", "时变信道、流量预测、连续观测建模。"),
        (7, "LSTM：增强长期记忆", "在 RNN 基础上加入细胞状态和门控。",
         PURPLE, "遗忘门、输入门、输出门共同决定哪些信息保留、写入和输出。",
         [("遗忘门", PURPLE), ("细胞状态", BLUE), ("输入门", PURPLE), ("输出", GREEN)],
         "较长序列、长期依赖任务。", "记忆更强，但参数更多、计算更复杂。", "连续 CSI 恢复、长期流量趋势预测。"),
        (8, "GRU：更简化的门控循环网络", "用更少的门控结构控制记忆。",
         GREEN, "相比 LSTM 更轻量，通常训练更快，适合对效率有要求的序列任务。",
         [("更新门", GREEN), ("重置门", GREEN), ("状态", BLUE)],
         "中等长度序列、轻量序列建模。", "比 LSTM 简洁，表达能力通常略弱一些。", "轻量 CSI 预测、实时状态估计。"),
        (9, "GAN：通过对抗生成数据", "生成器和判别器相互博弈。",
         AMBER, "生成器负责造样本，判别器负责辨真假，训练目标是让生成样本越来越像真实分布。",
         [("噪声", AMBER), ("生成器", AMBER), ("样本", GREEN), ("判别器", PURPLE)],
         "图像生成、数据增强、样本生成。", "目标是生成新样本，训练不稳定。", "信道样本生成、稀缺场景数据增强。"),
        (10, "Autoencoder：压缩与重建", "编码器压缩，解码器重建。",
         PURPLE, "模型通过重建输入来学习数据表示，关键是压缩后还能保留重要信息。",
         [("输入", BLUE), ("编码", PURPLE), ("低维码字", CYAN), ("重建", GREEN)],
         "降维、去噪、压缩、特征学习。", "是协作重建，不是对抗生成。", "CSI 反馈压缩、端到端物理层理解。"),
        (11, "Transformer：用注意力建模全局关系", "Self-Attention 直接计算位置间相关性。",
         RED, "不再像 RNN 一步步递推，而是让每个位置直接关注其他位置，并用位置编码补充顺序信息。",
         [("Query", RED), ("Key", AMBER), ("Value", GREEN), ("Attention", PURPLE)],
         "长序列、大模型、多模态、全局依赖。", "并行能力强，但计算和显存开销较大。", "多用户、多天线、多子载波关系建模。"),
    ]
    for spec in model_specs:
        n, title, subtitle, accent, core, diagram, fit, diff, comm = spec
        add_model_slide(prs, blank, n, title, subtitle, accent, core, diagram, fit, diff, comm)
        slides.append(("model", title, subtitle, False))

    # 12
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "八类模型区别总结", "不是谁绝对更强，而是面对的数据结构和任务目标不同。", accent=BLUE)
    rows = [
        ("模型", "核心特点", "适合任务", "主要局限"),
        ("FNN", "固定维度映射", "分类、回归、检测", "不利用空间/时间结构"),
        ("CNN", "局部感受野、权值共享", "图像、频谱图、矩阵", "全局依赖需多层组合"),
        ("RNN", "隐藏状态递推", "序列、时间变化", "长序列训练困难"),
        ("LSTM", "细胞状态和门控", "长期依赖", "计算复杂"),
        ("GRU", "简化门控", "轻量序列任务", "表达略弱于 LSTM"),
        ("GAN", "生成器 vs 判别器", "生成、数据增强", "训练不稳定"),
        ("Autoencoder", "编码压缩、解码重建", "降维、去噪、CSI反馈", "可能只记住输入"),
        ("Transformer", "自注意力、全局关系", "长序列、大模型、多模态", "算力和显存开销大"),
    ]
    add_table(s, 0.65, 1.95, [1.35, 3.25, 3.25, 3.45], rows, row_h=0.47, font_size=10)
    add_footer(s, 12)
    slides.append(("table", "八类模型区别总结", "", False))

    # 13
    s = prs.slides.add_slide(blank)
    add_bg(s, SLATE)
    add_title(s, "为什么通信任务能和深度学习结合", "很多通信问题可以翻译成学习问题。", dark=True, accent=AMBER)
    flows = [
        ("接收信号", "发送符号", "检测", BLUE),
        ("高维 CSI", "低维反馈/重建", "压缩", PURPLE),
        ("信道状态", "功率/波束", "优化映射", GREEN),
        ("真实/仿真信道", "更多样本", "生成", AMBER),
    ]
    for i, (a, b, c, col) in enumerate(flows):
        y = 2.25 + i * 0.85
        add_node(s, a, 1.15, y, 1.55, 0.45, col, 12)
        add_line(s, 2.85, y + 0.23, 3.55, y + 0.23, "94A3B8", 1.5, True)
        add_node(s, b, 3.75, y, 2.05, 0.45, col, 12)
        add_text(s, c, 6.25, y + 0.06, 2.3, 0.28, 15, "E5E7EB", True)
    add_text(s, "先判断通信问题更像映射、检测、压缩、预测还是生成，再选择网络结构。", 8.5, 2.7, 3.6, 1.0, 23, WHITE, True, line_spacing=1.08)
    add_footer(s, 13, dark=True)
    slides.append(("comm", "为什么通信任务能和深度学习结合", "", True))

    # 14
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "通信结合方向一：信道估计与信号检测", "对应 FNN / DNN，也可以结合模型驱动网络。", accent=BLUE)
    add_text(s, "基本形式", 0.85, 2.05, 1.2, 0.28, 13, BLUE, True)
    add_node(s, "接收信号 / 导频", 1.0, 2.55, 1.55, 0.48, BLUE, 11)
    add_line(s, 2.65, 2.79, 3.25, 2.79, "94A3B8", 1.6, True)
    add_node(s, "DNN", 3.45, 2.55, 1.0, 0.48, GREEN, 12)
    add_line(s, 4.55, 2.79, 5.15, 2.79, "94A3B8", 1.6, True)
    add_node(s, "信道状态 / 发送符号", 5.35, 2.55, 1.9, 0.48, AMBER, 11)
    add_text(s, "代表工作", 8.05, 2.05, 1.2, 0.28, 13, AMBER, True)
    add_text(s, "Ye、Li、Juang 2017/2018", 8.05, 2.48, 4.25, 0.34, 18, INK, True)
    add_text(s, "OFDM 信道估计与信号检测", 8.05, 2.95, 4.25, 0.34, 18, INK, True)
    add_text(s, "Power of Deep Learning for Channel Estimation and Signal Detection in OFDM Systems", 8.05, 3.35, 4.35, 0.52, 11, MUTED, True)
    add_text(s, "理解重点：不是简单抛弃 LS/MMSE，而是在导频少、无 CP、非线性失真等条件下，让模型学习接收信号中的规律。", 0.95, 4.55, 10.8, 0.82, 23, INK, True)
    add_footer(s, 14)
    slides.append(("case", "通信结合方向一：信道估计与信号检测", "", False))

    # 15
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "通信结合方向二：CSI 反馈压缩", "对应 Autoencoder + CNN/LSTM。", accent=PURPLE)
    chain = [("高维 CSI", BLUE), ("编码器\nUE 端", PURPLE), ("低维码字\n反馈", CYAN), ("解码器\nBS 端", PURPLE), ("恢复 CSI", GREEN)]
    for i, (txt, col) in enumerate(chain):
        x = 0.85 + i * 2.1
        add_node(s, txt, x, 3.05, 1.35, 0.65, col, 11)
        if i < len(chain) - 1:
            add_line(s, x + 1.42, 3.38, x + 1.95, 3.38, "94A3B8", 1.5, True)
    add_text(s, "CsiNet 可以理解为把 CSI 反馈问题写成压缩与重建：用户端压缩 CSI，基站端根据反馈码字恢复 CSI。", 0.95, 4.75, 5.4, 0.8, 20, INK, True)
    add_text(s, "进一步：CNN 利用空间结构；LSTM / GRU 利用连续时刻 CSI 的时间相关性。", 7.0, 4.75, 5.0, 0.8, 20, INK, True)
    add_footer(s, 15)
    slides.append(("case", "通信结合方向二：CSI 反馈压缩", "", False))

    # 16
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "通信结合方向三：调制识别 / 信号识别", "对应 CNN / RNN / Transformer。", accent=CYAN)
    add_text(s, "输入可以是 IQ 信号、频谱图或时频图。关键不是模型名字，而是这些输入里面确实存在可学习的局部模式和时间变化。", 0.88, 2.05, 10.9, 0.72, 24, INK, True)
    parts = [("IQ 片段", "局部波形模式", CNN, CYAN) if False else ("IQ 片段", "局部波形模式", "CNN", CYAN),
             ("连续观测", "时间变化", "RNN / LSTM / GRU", GREEN),
             ("长距离关系", "全局依赖", "Transformer", RED)]
    for i, (a, b, c, col) in enumerate(parts):
        x = 1.05 + i * 4.0
        add_pill(s, a, x, 3.55, 1.25, 0.38, col, size=11)
        add_text(s, b, x, 4.14, 2.5, 0.28, 16, INK, True)
        add_text(s, c, x, 4.62, 2.7, 0.33, 18, col, True)
    add_footer(s, 16)
    slides.append(("case", "通信结合方向三：调制识别 / 信号识别", "", False))

    # 17
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "通信结合方向四：资源分配与波束成形", "对应 DNN / 模型驱动深度学习。", accent=GREEN)
    add_text(s, "传统方法", 0.95, 2.08, 1.3, 0.28, 14, BLUE, True)
    add_node(s, "信道状态", 1.0, 2.65, 1.25, 0.45, BLUE, 11)
    add_line(s, 2.35, 2.88, 2.95, 2.88, "94A3B8", 1.5, True)
    add_node(s, "WMMSE 等\n迭代优化", 3.15, 2.55, 1.45, 0.62, BLUE, 11)
    add_line(s, 4.75, 2.88, 5.35, 2.88, "94A3B8", 1.5, True)
    add_node(s, "功率 / 波束", 5.55, 2.65, 1.25, 0.45, GREEN, 11)
    add_text(s, "深度学习思路", 0.95, 4.38, 1.6, 0.28, 14, GREEN, True)
    add_text(s, "用 DNN 学习“信道状态 → 优化结果”的映射，减少在线迭代计算。", 1.0, 4.82, 5.6, 0.48, 22, INK, True)
    add_text(s, "需要注意", 8.0, 2.08, 1.3, 0.28, 14, AMBER, True)
    add_text(s, "约束是否满足？\n换信道环境后是否泛化？\nUE/边缘设备复杂度能否接受？", 8.0, 2.55, 4.1, 1.35, 20, INK, True, line_spacing=1.22)
    add_footer(s, 17)
    slides.append(("case", "通信结合方向四：资源分配与波束成形", "", False))

    # 18
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "通信结合方向五：信道建模与数据增强", "对应 GAN。", accent=AMBER)
    add_text(s, "真实测量成本高，传统统计模型又未必覆盖复杂场景。GAN 可以尝试学习信道数据分布，再生成新的样本。", 0.9, 2.0, 10.7, 0.72, 24, INK, True)
    add_node(s, "真实/仿真信道", 1.0, 3.35, 1.6, 0.5, BLUE, 11)
    add_line(s, 2.75, 3.60, 3.35, 3.60, "94A3B8", 1.5, True)
    add_node(s, "GAN 学分布", 3.55, 3.35, 1.45, 0.5, AMBER, 11)
    add_line(s, 5.15, 3.60, 5.75, 3.60, "94A3B8", 1.5, True)
    add_node(s, "生成信道样本", 5.95, 3.35, 1.65, 0.5, GREEN, 11)
    add_text(s, "谨慎点：生成样本符合统计特征，不等于它在所有算法评估中都可靠。", 8.0, 3.05, 3.9, 0.85, 21, RED, True)
    add_footer(s, 18)
    slides.append(("case", "通信结合方向五：信道建模与数据增强", "", False))

    # 19
    s = prs.slides.add_slide(blank)
    add_bg(s, SLATE)
    add_title(s, "目前的阶段性理解", "我现在更倾向于把深度学习看作通信算法的补充工具。", dark=True, accent=GREEN)
    points = [
        ("不是替代", "不是直接替代通信理论，而是补充难建模、难优化、难实时的部分。", CYAN),
        ("先看问题", "先判断任务像映射、检测、压缩、预测还是生成，再选网络。", AMBER),
        ("更稳方向", "模型驱动 + 数据驱动，比纯黑箱更适合通信算法学习。", GREEN),
    ]
    for i, (h, b, col) in enumerate(points):
        y = 2.05 + i * 1.18
        add_pill(s, h, 0.95, y, 1.15, 0.38, col, size=11)
        add_text(s, b, 2.45, y - 0.02, 8.6, 0.52, 20, WHITE, True)
    add_text(s, "后续继续看：CsiNet 系列、Deep Unfolding、资源分配、GNN / Transformer 在通信中的应用。", 0.98, 6.12, 10.8, 0.4, 18, "D1FAE5", True)
    add_footer(s, 19, dark=True)
    slides.append(("end", "目前的阶段性理解", "", True))

    prs.save(PPTX_PATH)
    return slides


def draw_text(draw, text, xy, size=28, color=INK, bold=False, max_width=None, anchor=None):
    font = font_px(size, bold)
    fill = "#" + color.strip("#")
    if max_width:
        x, y = xy
        for line in wrap_text(draw, text, font, max_width):
            draw.text((x, y), line, font=font, fill=fill)
            box = draw.textbbox((0, 0), "深", font=font)
            y += int((box[3] - box[1]) * 1.22)
        return
    draw.text(xy, text, font=font, fill=fill, anchor=anchor)


def wrap_text(draw, text, font, max_width):
    lines = []
    for para in text.split("\n"):
        current = ""
        for ch in para:
            trial = current + ch
            if draw.textbbox((0, 0), trial, font=font)[2] <= max_width or not current:
                current = trial
            else:
                lines.append(current)
                current = ch
        if current:
            lines.append(current)
    return lines


def emu_to_px(value: int, axis: str) -> int:
    slide_emu = int((SLIDE_W if axis == "x" else SLIDE_H) * 914400)
    slide_px = PX_W if axis == "x" else PX_H
    return int(value / slide_emu * slide_px)


def first_srgb(node, ns):
    color = node.find(".//a:srgbClr", ns)
    if color is not None:
        return color.attrib.get("val")
    return None


def extract_paragraphs(tx_body, ns):
    paras = []
    for para in tx_body.findall("a:p", ns):
        text = "".join(t.text or "" for t in para.findall(".//a:t", ns))
        if text.strip():
            paras.append(text)
    return "\n".join(paras)


def text_style(tx_body, ns, fallback=INK):
    rpr = tx_body.find(".//a:rPr", ns)
    if rpr is None:
        return 14, fallback, False
    size = int(rpr.attrib.get("sz", "1400")) // 100
    color = first_srgb(rpr, ns) or fallback
    bold = rpr.attrib.get("b") == "1"
    return max(7, size), color, bold


def text_align(tx_body, ns):
    ppr = tx_body.find(".//a:pPr", ns)
    if ppr is None:
        return "left"
    return "center" if ppr.attrib.get("algn") == "ctr" else "left"


def draw_text_box(draw, text, box, size, color, bold=False, align="left"):
    x, y, w, h = box
    pad = 8
    size = max(7, min(size, 56))
    while size > 7:
        font = font_px(size, bold)
        lines = []
        for para in text.split("\n"):
            lines.extend(wrap_text(draw, para, font, max(20, w - pad * 2)))
        sample = draw.textbbox((0, 0), "深", font=font)
        line_h = max(12, int((sample[3] - sample[1]) * 1.15))
        if len(lines) * line_h <= h - pad or size <= 9:
            break
        size -= 1
    yy = y + pad
    if h < 90 or align == "center":
        yy = y + max(pad, int((h - len(lines) * line_h) / 2))
    for line in lines:
        if align == "center":
            draw.text((x + w / 2, yy), line, font=font, fill="#" + color, anchor="ma")
        else:
            draw.text((x + pad, yy), line, font=font, fill="#" + color)
        yy += line_h


def render_pptx_previews(path: Path):
    ns = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }
    with ZipFile(path) as zf:
        slide_names = sorted(
            [n for n in zf.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", n)],
            key=lambda n: int(re.search(r"slide(\d+)\.xml", n).group(1)),
        )
        for idx, name in enumerate(slide_names, start=1):
            root = ET.fromstring(zf.read(name))
            bg_node = root.find(".//p:bg", ns)
            bg_color = first_srgb(bg_node if bg_node is not None else root, ns) or LIGHT
            im = Image.new("RGB", (PX_W, PX_H), "#" + bg_color)
            draw = ImageDraw.Draw(im)
            for sp in root.findall(".//p:sp", ns):
                xfrm = sp.find(".//p:spPr/a:xfrm", ns)
                if xfrm is None:
                    continue
                off = xfrm.find("a:off", ns)
                ext = xfrm.find("a:ext", ns)
                if off is None or ext is None:
                    continue
                x = emu_to_px(int(off.attrib.get("x", "0")), "x")
                y = emu_to_px(int(off.attrib.get("y", "0")), "y")
                w = emu_to_px(int(ext.attrib.get("cx", "0")), "x")
                h = emu_to_px(int(ext.attrib.get("cy", "0")), "y")
                if w <= 0 or h <= 0:
                    continue

                sppr = sp.find("p:spPr", ns)
                fill = first_srgb(sppr, ns) if sppr is not None else None
                geom = sp.find(".//a:prstGeom", ns)
                geom_name = geom.attrib.get("prst", "rect") if geom is not None else "rect"
                if fill:
                    if geom_name in {"ellipse"}:
                        draw.ellipse((x, y, x + w, y + h), fill="#" + fill)
                    elif geom_name in {"roundRect"}:
                        draw.rounded_rectangle((x, y, x + w, y + h), radius=max(6, min(w, h) // 5), fill="#" + fill)
                    else:
                        draw.rectangle((x, y, x + w, y + h), fill="#" + fill)

                tx_body = sp.find("p:txBody", ns)
                if tx_body is not None:
                    text = extract_paragraphs(tx_body, ns)
                    if text:
                        fallback = WHITE if fill and fill.upper() in {SLATE, BLUE, CYAN, GREEN, AMBER, PURPLE, RED} else INK
                        size, color, bold = text_style(tx_body, ns, fallback=fallback)
                        draw_text_box(draw, text, (x, y, w, h), size, color, bold, text_align(tx_body, ns))
            im.save(PREVIEW_DIR / f"slide_{idx:02d}.png")


def make_previews(slides):
    for old in PREVIEW_DIR.glob("slide_*.png"):
        old.unlink()
    render_pptx_previews(PPTX_PATH)
    thumbs = [Image.open(p).resize((480, 270)) for p in sorted(PREVIEW_DIR.glob("slide_*.png"))]
    cols = 4
    rows = math.ceil(len(thumbs) / cols)
    montage = Image.new("RGB", (cols * 480, rows * 300), "#F8FAFC")
    d = ImageDraw.Draw(montage)
    for i, im in enumerate(thumbs):
        x = (i % cols) * 480
        y = (i // cols) * 300
        montage.paste(im, (x, y))
        d.text((x + 12, y + 274), f"Slide {i + 1:02d}", fill="#475569", font=font_px(14, True))
    montage.save(MONTAGE_PATH)


def inspect_pptx(path: Path):
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    issues = []
    with ZipFile(path) as zf:
        slide_names = sorted(
            [n for n in zf.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", n)],
            key=lambda n: int(re.search(r"slide(\d+)\.xml", n).group(1)),
        )
        for idx, name in enumerate(slide_names, start=1):
            root = ET.fromstring(zf.read(name))
            texts = [t.text or "" for t in root.findall(".//a:t", ns)]
            joined = "".join(texts)
            if "Click to add" in joined or "Slide Number" in joined:
                issues.append(f"slide {idx}: placeholder text found")
            if not joined.strip():
                issues.append(f"slide {idx}: no text found")
    return len(slide_names), issues


def main():
    slides = create_deck()
    make_previews(slides)
    count, issues = inspect_pptx(PPTX_PATH)
    print(f"created={PPTX_PATH}")
    print(f"slides={count}")
    print(f"previews={PREVIEW_DIR}")
    print(f"montage={MONTAGE_PATH}")
    if issues:
        for issue in issues:
            print(f"- {issue}")
        raise SystemExit(1)
    print("pptx_xml_inspection=ok")


if __name__ == "__main__":
    main()
